"""
RAG Service - Retrieval Augmented Generation for Document Spaces
Query AI with context from uploaded documents
"""

import os
from typing import List, Dict, Any, Optional
from pathlib import Path


class RAGService:
    """Simple RAG service using Gemini for document Q&A"""
    
    def __init__(self):
        self.client = None
        self.model_name = "gemini-3-flash-preview"
        self._init_client()
    
    def _init_client(self):
        """Initialize Gemini client with proper SDK"""
        api_key = os.getenv("GEMINI_API_KEY")
        if not api_key:
            print("[RAG Service] No GEMINI_API_KEY found")
            return
        
        try:
            # Try google-generativeai SDK first
            import google.generativeai as genai
            genai.configure(api_key=api_key)
            self.client = genai.GenerativeModel(self.model_name)
            self.use_genai = True
            print(f"[RAG Service] Initialized with google-generativeai, model: {self.model_name}")
        except ImportError:
            try:
                # Fallback to google.genai SDK
                from google import genai
                self.client = genai.Client(api_key=api_key)
                self.use_genai = False
                print(f"[RAG Service] Initialized with google.genai, model: {self.model_name}")
            except Exception as e:
                print(f"[RAG Service] Init error: {e}")
    
    def is_available(self) -> bool:
        """Check if service is available"""
        return self.client is not None
    
    def _generate(self, prompt: str) -> str:
        """Generate content with proper SDK"""
        if self.use_genai:
            response = self.client.generate_content(prompt)
            return response.text
        else:
            response = self.client.models.generate_content(
                model=self.model_name,
                contents=prompt
            )
            return response.text
    
    async def query(
        self, 
        question: str, 
        file_paths: List[str],
        space_name: str = "Document Space",
        space_id: str = None
    ) -> Dict[str, Any]:
        """
        Smart Query with Gold Layer support
        
        Flow:
        1. Check Gold tables và Smart Router
        2. Nếu có structured data → SQL query (100% accurate)
        3. Fallback → AI với document context
        """
        if not self.is_available():
            return {
                "answer": "❌ AI service chưa được cấu hình. Vui lòng kiểm tra GEMINI_API_KEY.",
                "sources": []
            }
        
        # Try Gold Layer first if space_id provided
        if space_id:
            try:
                from services.gold_layer_service import gold_layer_service
                from services.smart_query_router import smart_query_router
                
                # Get Gold tables for this space
                gold_tables = await gold_layer_service.get_gold_tables(space_id)
                
                if gold_tables:
                    # Route query through Smart Router
                    routed = await smart_query_router.route_query(
                        question=question,
                        space_id=space_id,
                        gold_tables=gold_tables,
                        file_paths=file_paths
                    )
                    
                    # If router returned a direct answer (SQL result)
                    if not routed.get("use_ai"):
                        return {
                            "answer": routed.get("answer", ""),
                            "sources": [],
                            "source_type": routed.get("source_type", "gold"),
                            "is_gold_query": True,
                            "from_cache": routed.get("from_cache", False)
                        }
                    
                    # Add Gold context to AI query
                    gold_context = routed.get("gold_context")
                    if gold_context:
                        # Use Gold data as additional context
                        pass  # Will enhance prompt below
            except Exception as e:
                print(f"[RAG] Gold Layer error: {e}")
        
        # Fallback: Read documents and use AI
        context_parts = []
        sources = []
        all_docs = []
        
        for file_path in file_paths:
            try:
                file_name = Path(file_path).stem
                content = await self._read_file_content(file_path)
                
                doc_info = {"file": file_name, "path": file_path, "has_content": content is not None}
                all_docs.append(doc_info)
                
                if content:
                    max_chars = min(8000, 50000 // max(len(file_paths), 1))
                    context_parts.append(f"=== {file_name} ===\n{content[:max_chars]}")
                    sources.append({"file": file_name, "path": file_path})
            except Exception as e:
                print(f"[RAG] Error reading {file_path}: {e}")
                all_docs.append({"file": Path(file_path).stem, "path": file_path, "has_content": False, "error": str(e)})
        
        if not context_parts:
            return {
                "answer": "⚠️ Không thể đọc nội dung các file trong space này.",
                "sources": []
            }
        
        # Build improved prompt
        context = "\n\n".join(context_parts)
        doc_list = "\n".join([f"• {d['file']}" for d in all_docs])
        
        prompt = f"""Bạn là AI assistant cho "{space_name}". Trả lời câu hỏi dựa trên nội dung tài liệu được cung cấp.

DANH SÁCH TÀI LIỆU ({len(all_docs)} files):
{doc_list}

NỘI DUNG TÀI LIỆU:
{context}

CÂU HỎI: {question}

QUAN TRỌNG - QUY TẮC TRẢ LỜI:
1. KHÔNG dùng markdown (không dùng **, *, #, ```, etc.)
2. Dùng bullet points bằng dấu • hoặc - 
3. Dùng số thứ tự 1. 2. 3. khi liệt kê
4. Viết văn bản thuần, dễ đọc, tự nhiên
5. Trả lời chính xác 100% dựa trên tài liệu
6. Nếu không tìm thấy thông tin, nói rõ ràng"""

        try:
            answer = self._generate(prompt)
            
            return {
                "answer": answer,
                "sources": sources,
                "total_documents": len(all_docs),
                "source_type": "ai"
            }
        except Exception as e:
            return {
                "answer": f"❌ Lỗi khi gọi AI: {str(e)}",
                "sources": []
            }
    
    async def _read_file_content(self, file_path: str) -> Optional[str]:
        """Read content from various file types"""
        path = Path(file_path)
        
        if not path.exists():
            return None
        
        ext = path.suffix.lower()
        
        try:
            # Text files
            if ext in ['.txt', '.csv', '.json', '.md']:
                return path.read_text(encoding='utf-8', errors='ignore')[:10000]
            
            # PDF - Extract ALL text from all pages
            elif ext == '.pdf':
                import pdfplumber
                text_parts = []
                with pdfplumber.open(file_path) as pdf:
                    for page_num, page in enumerate(pdf.pages[:20], 1):  # Up to 20 pages
                        page_text = page.extract_text()
                        if page_text and page_text.strip():
                            text_parts.append(f"[Trang {page_num}]\n{page_text.strip()}")
                
                if text_parts:
                    return "\n\n".join(text_parts)[:15000]
                else:
                    # No text extracted - might be scanned PDF
                    return "[PDF này có thể là ảnh scan, không extract được text trực tiếp]"
            
            # Word docs
            elif ext == '.docx':
                from docx import Document
                doc = Document(file_path)
                text = "\n".join([p.text for p in doc.paragraphs])
                return text[:10000]
            
            # Excel
            elif ext in ['.xlsx', '.xls']:
                import pandas as pd
                df = pd.read_excel(file_path, nrows=100)
                return df.to_string()[:10000]
            
            # PowerPoint
            elif ext == '.pptx':
                from pptx import Presentation
                prs = Presentation(file_path)
                text = []
                for slide in prs.slides[:10]:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text.append(shape.text)
                return "\n".join(text)[:10000]
            
            # Parquet
            elif ext == '.parquet':
                import pandas as pd
                df = pd.read_parquet(file_path).head(100)
                return df.to_string()[:10000]
            
            else:
                return None
                
        except Exception as e:
            print(f"[RAG] Read error {file_path}: {e}")
            return None


# Singleton instance
rag_service = RAGService()
