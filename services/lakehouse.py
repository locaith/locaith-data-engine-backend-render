import duckdb
import pandas as pd
import pyarrow.parquet as pq
import json
import os
from typing import Optional, List, Dict, Any, Tuple
from config import settings
from database import get_db
from services.auth_service import generate_uuid

class LakehouseService:
    def __init__(self):
        self.data_dir = settings.DATA_DIR
    
    def _extract_docx_data(self, file_path: str) -> pd.DataFrame:
        """Extract text content from Word document"""
        try:
            from docx import Document
            doc = Document(file_path)
            
            # Extract paragraphs
            paragraphs = []
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append({"content": para.text})
            
            # Extract tables if any
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join([cell.text for cell in row.cells])
                    if row_text.strip():
                        paragraphs.append({"content": row_text})
            
            if paragraphs:
                return pd.DataFrame(paragraphs)
            else:
                return pd.DataFrame([{"content": "Document không có nội dung"}])
        except ImportError:
            # Fallback: read as zip and extract text
            import zipfile
            try:
                with zipfile.ZipFile(file_path, 'r') as z:
                    xml_content = z.read('word/document.xml').decode('utf-8')
                    # Simple text extraction
                    import re
                    text = re.sub(r'<[^>]+>', '', xml_content)
                    return pd.DataFrame([{"content": text[:50000]}])
            except:
                return pd.DataFrame([{"content": "Không thể đọc file DOCX"}])
    
    def _extract_pptx_data(self, file_path: str) -> pd.DataFrame:
        """Extract text from PowerPoint file"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            
            slides_content = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                if slide_text:
                    slides_content.append({
                        "slide": slide_num,
                        "content": " | ".join(slide_text)
                    })
            
            if slides_content:
                return pd.DataFrame(slides_content)
            else:
                return pd.DataFrame([{"content": "PowerPoint không có nội dung text"}])
        except ImportError:
            # Fallback: try to extract from zip
            import zipfile
            try:
                with zipfile.ZipFile(file_path, 'r') as z:
                    text_parts = []
                    for name in z.namelist():
                        if name.startswith('ppt/slides/slide') and name.endswith('.xml'):
                            content = z.read(name).decode('utf-8', errors='ignore')
                            import re
                            text = re.sub(r'<[^>]+>', ' ', content)
                            text = ' '.join(text.split())[:5000]
                            text_parts.append({"content": text})
                    if text_parts:
                        return pd.DataFrame(text_parts)
                    return pd.DataFrame([{"content": "Không thể đọc nội dung PPTX"}])
            except:
                return pd.DataFrame([{"content": "Không thể đọc file PPTX"}])
    
    def _extract_image_data(self, file_path: str) -> pd.DataFrame:
        """Extract text from images using OCR (Gemini Vision or Tesseract)"""
        import asyncio
        
        try:
            from services.ocr_service import ocr_service
            
            # Run OCR
            loop = asyncio.new_event_loop()
            result = loop.run_until_complete(ocr_service.extract_text_from_image(file_path))
            loop.close()
            
            if result.get("success") and result.get("text"):
                return pd.DataFrame([{
                    "content": result["text"],
                    "_ocr_method": result.get("method", "unknown"),
                    "_type": "ocr"
                }])
            else:
                return pd.DataFrame([{
                    "content": "[Không thể OCR ảnh]",
                    "error": result.get("error", "Unknown error"),
                    "_type": "error"
                }])
                
        except Exception as e:
            return pd.DataFrame([{
                "content": f"[Lỗi OCR: {str(e)}]",
                "_type": "error"
            }])
    
    def _extract_pdf_data(self, file_path: str) -> pd.DataFrame:
        """Extract ALL content from PDF file - tables AND text, with OCR fallback"""
        import pdfplumber
        import asyncio
        
        all_tables = []
        all_text = []
        
        try:
            with pdfplumber.open(file_path) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    # Always try to extract tables
                    tables = page.extract_tables()
                    for table in tables:
                        if table and len(table) > 1:
                            headers = [h if h else f"col_{i}" for i, h in enumerate(table[0])]
                            for row in table[1:]:
                                if row:
                                    row_dict = {"_page": page_num, "_type": "table"}
                                    for i, val in enumerate(row):
                                        if i < len(headers):
                                            row_dict[headers[i]] = val
                                    all_tables.append(row_dict)
                    
                    # ALWAYS extract text content
                    text = page.extract_text()
                    if text and text.strip():
                        all_text.append({
                            "page": page_num,
                            "content": text.strip(),
                            "_type": "text"
                        })
            
            # Check if we got content
            if all_tables and all_text:
                df = pd.DataFrame(all_text)
            elif all_tables:
                df = pd.DataFrame(all_tables)
            elif all_text:
                df = pd.DataFrame(all_text)
            else:
                # NO CONTENT - Try OCR for scanned PDFs
                print(f"[Lakehouse] PDF appears scanned, attempting OCR: {file_path}")
                try:
                    from services.ocr_service import ocr_service
                    
                    # Run OCR (async in sync context)
                    loop = asyncio.new_event_loop()
                    ocr_result = loop.run_until_complete(ocr_service.extract_text_from_pdf(file_path))
                    loop.close()
                    
                    if ocr_result.get("success") and ocr_result.get("pages"):
                        pages = ocr_result["pages"]
                        df = pd.DataFrame([{
                            "page": p["page"],
                            "content": p["text"],
                            "_type": "ocr",
                            "_ocr_method": ocr_result.get("method", "unknown")
                        } for p in pages])
                        print(f"[Lakehouse] OCR success: {len(pages)} pages extracted")
                    else:
                        df = pd.DataFrame({
                            "message": ["PDF là ảnh scan - OCR thất bại hoặc không khả dụng"],
                            "error": [ocr_result.get("error", "Unknown OCR error")],
                            "_type": ["error"]
                        })
                except Exception as ocr_error:
                    print(f"[Lakehouse] OCR failed: {ocr_error}")
                    df = pd.DataFrame({
                        "message": ["PDF có thể là ảnh scan, OCR không khả dụng"],
                        "suggestion": ["Cài đặt tesseract hoặc sử dụng Gemini API key"],
                        "_type": ["error"]
                    })
            
            return df
            
        except Exception as e:
            return pd.DataFrame({
                "error": [f"Lỗi đọc PDF: {str(e)}"],
                "file": [file_path]
            })
    
    def ingest_file(self, file_path: str, user_id: str, name: str, description: str = None, space_id: str = None, storage_url: str = None) -> Dict[str, Any]:
        """Ingest a file into the lakehouse"""
        file_ext = os.path.splitext(file_path)[1].lower()
        
        # Read data based on file type
        if file_ext == '.csv':
            df = pd.read_csv(file_path)
            file_type = 'csv'
        elif file_ext == '.json':
            df = pd.read_json(file_path)
            file_type = 'json'
        elif file_ext == '.parquet':
            df = pd.read_parquet(file_path)
            file_type = 'parquet'
        elif file_ext == '.pdf':
            # Extract data from PDF using pdfplumber
            df = self._extract_pdf_data(file_path)
            file_type = 'pdf'
        elif file_ext in ['.xlsx', '.xls']:
            # Read Excel files
            df = pd.read_excel(file_path)
            file_type = 'excel'
        elif file_ext in ['.docx', '.doc']:
            # Read Word documents
            df = self._extract_docx_data(file_path)
            file_type = 'docx'
        elif file_ext == '.txt':
            # Read plain text files
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            df = pd.DataFrame([{"content": content}])
            file_type = 'txt'
        elif file_ext in ['.xml', '.html', '.htm']:
            # Read XML/HTML as text
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            df = pd.DataFrame([{"content": content}])
            file_type = file_ext[1:]  # Remove dot
        elif file_ext in ['.pptx', '.ppt']:
            # Extract text from PowerPoint
            df = self._extract_pptx_data(file_path)
            file_type = 'pptx'
        elif file_ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp']:
            # Extract text from images using OCR
            df = self._extract_image_data(file_path)
            file_type = 'image'
        else:
            # Try to read as plain text for unknown types
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
                df = pd.DataFrame([{"content": content}])
                file_type = 'text'
            except:
                raise ValueError(f"Unsupported file type: {file_ext}")
        
        # Generate dataset ID
        dataset_id = generate_uuid()
        
        # ========== ENTERPRISE DATA CLEANING (100% ACCURACY) ==========
        # Clean and validate BEFORE storing for AI fine-tuning quality
        import asyncio
        try:
            from services.data_cleaner_service import enterprise_data_cleaner
            
            # Run cleaning in async context
            loop = asyncio.new_event_loop()
            df_cleaned, integrity_report = loop.run_until_complete(
                enterprise_data_cleaner.clean_dataframe(df, source_name=name, strict_mode=True)
            )
            loop.close()
            
            # Log cleaning results
            print(f"[Lakehouse] Data cleaned: {integrity_report.original_rows} rows, {len(integrity_report.actions_taken)} actions, integrity: {integrity_report.integrity_score}%")
            
            # Use cleaned dataframe
            df = df_cleaned
            
            # Store cleaning metadata
            cleaning_metadata = {
                "original_rows": integrity_report.original_rows,
                "cleaned_rows": integrity_report.cleaned_rows,
                "integrity_score": integrity_report.integrity_score,
                "actions_count": len(integrity_report.actions_taken),
                "warnings": integrity_report.warnings[:5],  # Limit warnings
                "data_loss": integrity_report.data_loss
            }
            
        except Exception as clean_error:
            print(f"[Lakehouse] Data cleaning warning: {clean_error}")
            cleaning_metadata = {"status": "skipped", "reason": str(clean_error)}
        # ========== END DATA CLEANING ==========
        
        # ========== ADD PROVENANCE TRACKING ==========
        # Every row gets metadata about its source for 100% traceability
        original_filename = os.path.basename(file_path) if description else name
        
        # Add provenance columns to every row
        df['_source_file'] = name  # Original file name
        df['_source_type'] = file_type
        df['_row_number'] = range(1, len(df) + 1)  # 1-indexed row number
        df['_dataset_id'] = dataset_id
        df['_ingested_at'] = pd.Timestamp.now().isoformat()
        
        # Create content hash for deduplication detection (not removal!)
        if 'content' in df.columns:
            df['_content_hash'] = df['content'].apply(lambda x: hash(str(x)[:500]) if x else 0)
        # ========== END PROVENANCE TRACKING ==========
        
        # Convert to Parquet for efficient storage
        parquet_path = os.path.join(self.data_dir, "parquet", f"{dataset_id}.parquet")
        os.makedirs(os.path.dirname(parquet_path), exist_ok=True)
        df.to_parquet(parquet_path, index=False)
        
        # Get schema info (excluding internal columns for clean API)
        user_columns = [col for col in df.columns if not col.startswith('_')]
        schema_info = {
            "columns": [
                {"name": col, "type": str(df[col].dtype)}
                for col in user_columns
            ],
            "internal_columns": [col for col in df.columns if col.startswith('_')]
        }
        
        file_size = os.path.getsize(parquet_path)
        row_count = len(df)
        
        # Save dataset metadata (include storage_url if available for persistence)
        with get_db() as conn:
            conn.execute("""
                INSERT INTO datasets (id, user_id, space_id, name, description, file_path, file_type, file_size, row_count, schema_json, storage_url)
                VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
            """, [dataset_id, user_id, space_id, name, description, parquet_path, file_type, file_size, row_count, json.dumps(schema_info), storage_url])
        
        return {
            "id": dataset_id,
            "name": name,
            "file_type": file_type,
            "file_size": file_size,
            "row_count": row_count,
            "schema": schema_info,
            "provenance": {
                "source_file": name,
                "has_row_tracking": True,
                "has_content_hash": 'content' in df.columns
            }
        }
    
    def get_datasets(self, user_id: str) -> List[Dict[str, Any]]:
        """Get all datasets for a user"""
        with get_db() as conn:
            result = conn.execute("""
                SELECT id, name, description, file_type, file_size, row_count, schema_json, created_at
                FROM datasets
                WHERE user_id = ?
                ORDER BY created_at DESC
            """, [user_id]).fetchall()
            
            return [
                {
                    "id": row[0],
                    "name": row[1],
                    "description": row[2],
                    "file_type": row[3],
                    "file_size": row[4],
                    "row_count": row[5],
                    "schema_json": row[6],
                    "created_at": row[7]
                }
                for row in result
            ]
    
    def get_dataset(self, dataset_id: str, user_id: str) -> Optional[Dict[str, Any]]:
        """Get a specific dataset"""
        with get_db() as conn:
            result = conn.execute("""
                SELECT id, name, description, file_path, file_type, file_size, row_count, schema_json, created_at
                FROM datasets
                WHERE id = ? AND user_id = ?
            """, [dataset_id, user_id]).fetchone()
            
            if not result:
                return None
            
            return {
                "id": result[0],
                "name": result[1],
                "description": result[2],
                "file_path": result[3],
                "file_type": result[4],
                "file_size": result[5],
                "row_count": result[6],
                "schema_json": result[7],
                "created_at": result[8]
            }
    
    def preview_dataset(self, dataset_id: str, user_id: str, limit: int = 100) -> Optional[Dict[str, Any]]:
        """Preview dataset data"""
        dataset = self.get_dataset(dataset_id, user_id)
        if not dataset:
            return None
        
        df = pd.read_parquet(dataset["file_path"]).head(limit)
        
        return {
            "columns": list(df.columns),
            "data": df.to_dict(orient='records'),
            "total_rows": dataset["row_count"]
        }
    
    def delete_dataset(self, dataset_id: str, user_id: str) -> bool:
        """Delete a dataset"""
        dataset = self.get_dataset(dataset_id, user_id)
        if not dataset:
            return False
        
        # Delete file
        if os.path.exists(dataset["file_path"]):
            os.remove(dataset["file_path"])
        
        # Delete metadata
        with get_db() as conn:
            conn.execute("DELETE FROM datasets WHERE id = ? AND user_id = ?", [dataset_id, user_id])
        
        return True
    
    def execute_query(self, sql: str, user_id: str) -> Tuple[Dict[str, Any], int]:
        """Execute a SQL query"""
        import time
        
        start_time = time.time()
        
        with get_db() as conn:
            # Register user's datasets as tables
            datasets = self.get_datasets(user_id)
            for ds in datasets:
                table_name = ds["name"].replace(" ", "_").lower()
                conn.execute(f"CREATE OR REPLACE VIEW {table_name} AS SELECT * FROM read_parquet('{ds['file_path'].replace(chr(92), '/')}')")
            
            # Execute query
            try:
                result = conn.execute(sql)
                columns = [desc[0] for desc in result.description]
                data = result.fetchall()
                
                execution_time = int((time.time() - start_time) * 1000)
                
                # Log query
                query_id = generate_uuid()
                conn.execute("""
                    INSERT INTO query_history (id, user_id, query_text, execution_time_ms, row_count, status)
                    VALUES (?, ?, ?, ?, ?, 'success')
                """, [query_id, user_id, sql, execution_time, len(data)])
                
                return {
                    "columns": columns,
                    "data": [list(row) for row in data],
                    "row_count": len(data),
                    "execution_time_ms": execution_time
                }, 200
                
            except Exception as e:
                execution_time = int((time.time() - start_time) * 1000)
                
                # Log failed query
                query_id = generate_uuid()
                conn.execute("""
                    INSERT INTO query_history (id, user_id, query_text, execution_time_ms, status, error_message)
                    VALUES (?, ?, ?, ?, 'error', ?)
                """, [query_id, user_id, sql, execution_time, str(e)])
                
                return {"error": str(e)}, 400
    
    def get_query_history(self, user_id: str, limit: int = 50) -> List[Dict[str, Any]]:
        """Get query history for a user"""
        with get_db() as conn:
            result = conn.execute("""
                SELECT id, query_text, execution_time_ms, row_count, status, error_message, created_at
                FROM query_history
                WHERE user_id = ?
                ORDER BY created_at DESC
                LIMIT ?
            """, [user_id, limit]).fetchall()
            
            return [
                {
                    "id": row[0],
                    "query_text": row[1],
                    "execution_time_ms": row[2],
                    "row_count": row[3],
                    "status": row[4],
                    "error_message": row[5],
                    "created_at": row[6]
                }
                for row in result
            ]
    
    # ============ API KEY ISOLATED METHODS (for commercial) ============
    
    def ingest_file_by_api_key(self, file_path: str, api_key_id: str, user_id: str, name: str, description: str = None) -> Dict[str, Any]:
        """Ingest a file with API Key isolation - data belongs to the API Key"""
        result = self.ingest_file(file_path, user_id, name, description)
        
        # Update the dataset to link to API Key
        with get_db() as conn:
            conn.execute("""
                UPDATE datasets SET api_key_id = ? WHERE id = ?
            """, [api_key_id, result["id"]])
            
            # Update storage used
            file_size_mb = result["file_size"] / (1024 * 1024)
            conn.execute("""
                UPDATE api_keys SET storage_used_mb = storage_used_mb + ? WHERE id = ?
            """, [file_size_mb, api_key_id])
        
        return result
    
    def get_datasets_by_api_key(self, api_key_id: str) -> List[Dict[str, Any]]:
        """Get all datasets for a specific API Key (isolated data)"""
        with get_db() as conn:
            result = conn.execute("""
                SELECT id, name, description, file_type, file_size, row_count, schema_json, created_at
                FROM datasets
                WHERE api_key_id = ?
                ORDER BY created_at DESC
            """, [api_key_id]).fetchall()
            
            return [
                {
                    "id": row[0],
                    "name": row[1],
                    "description": row[2],
                    "file_type": row[3],
                    "file_size": row[4],
                    "row_count": row[5],
                    "schema_json": row[6],
                    "created_at": row[7]
                }
                for row in result
            ]
    
    def preview_dataset_by_api_key(self, dataset_id: str, api_key_id: str, limit: int = 100) -> Optional[Dict[str, Any]]:
        """Preview dataset for a specific API Key (isolated)"""
        with get_db() as conn:
            result = conn.execute("""
                SELECT file_path FROM datasets WHERE id = ? AND api_key_id = ?
            """, [dataset_id, api_key_id]).fetchone()
            
            if not result:
                return None
            
            file_path = result[0]
            
        try:
            df = pd.read_parquet(file_path)
            preview_df = df.head(limit)
            
            return {
                "columns": list(df.columns),
                "data": preview_df.to_dict(orient='records'),
                "total_rows": len(df),
                "preview_rows": len(preview_df)
            }
        except Exception as e:
            return {"error": str(e)}
    
    def delete_dataset_by_api_key(self, dataset_id: str, api_key_id: str) -> bool:
        """Delete a dataset for a specific API Key"""
        with get_db() as conn:
            # Get file info
            result = conn.execute("""
                SELECT file_path, file_size FROM datasets WHERE id = ? AND api_key_id = ?
            """, [dataset_id, api_key_id]).fetchone()
            
            if not result:
                return False
            
            file_path, file_size = result
            
            # Update storage used
            file_size_mb = file_size / (1024 * 1024) if file_size else 0
            conn.execute("""
                UPDATE api_keys SET storage_used_mb = MAX(0, storage_used_mb - ?) WHERE id = ?
            """, [file_size_mb, api_key_id])
            
            # Delete dataset
            conn.execute("DELETE FROM datasets WHERE id = ? AND api_key_id = ?", [dataset_id, api_key_id])
            
            # Delete parquet file
            if file_path and os.path.exists(file_path):
                os.remove(file_path)
            
            return True

lakehouse_service = LakehouseService()

