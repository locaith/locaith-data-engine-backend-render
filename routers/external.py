"""
External Data API - Returns RAW data for 3rd party AI to process
Instead of processing with internal AI, this endpoint returns document content
so external AI services can use it as context
"""

from fastapi import APIRouter, HTTPException, Header, Request
from typing import Optional, List
from pydantic import BaseModel
from pathlib import Path

from services.api_key_service import api_key_service
from database import get_db

router = APIRouter(prefix="/external", tags=["External API - Raw Data"])


class DataQueryRequest(BaseModel):
    query: Optional[str] = None  # Optional search query to filter results
    limit: int = 5  # Max number of documents to return


class DocumentContent(BaseModel):
    file_name: str
    file_type: str
    content: str  # Raw text content
    

class DataQueryResponse(BaseModel):
    space_name: str
    documents: List[DocumentContent]
    total_documents: int


async def validate_api_key(x_api_key: str = Header(None, alias="X-API-Key")):
    """Validate API Key and return key info"""
    if not x_api_key:
        raise HTTPException(
            status_code=401,
            detail="API Key không được cung cấp. Vui lòng thêm header X-API-Key"
        )
    
    key_info = api_key_service.validate_key(x_api_key)
    if not key_info:
        raise HTTPException(
            status_code=401,
            detail="API Key không hợp lệ hoặc đã hết hạn"
        )
    
    return key_info


@router.post("/data", response_model=DataQueryResponse)
async def get_raw_data(
    request: DataQueryRequest,
    x_api_key: str = Header(..., alias="X-API-Key")
):
    """
    Get RAW document content for 3rd party AI processing.
    
    This endpoint returns the actual content of documents in the space
    linked to the API key, instead of processing with internal AI.
    
    Use case: Your AI chatbot calls this API, gets raw data, and uses 
    your own AI (Gemini, GPT, Claude) to generate answers.
    """
    try:
        print(f"[External API] Validating API key: {x_api_key[:10]}...")
        key_info = await validate_api_key(x_api_key)
        print(f"[External API] Key info: {key_info}")
        
        space_id = key_info.get("space_id")
        if not space_id:
            raise HTTPException(
                status_code=400,
                detail="API Key này chưa được gắn với Document Space nào"
            )
        
        print(f"[External API] Getting space: {space_id}")
        # Get space info
        with get_db() as conn:
            space = conn.execute("""
                SELECT id, name FROM document_spaces WHERE id = ?
            """, [space_id]).fetchone()
            
            if not space:
                raise HTTPException(
                    status_code=404,
                    detail="Document Space không tồn tại"
                )
            
            print(f"[External API] Space found: {space[1]}")
            
            # Get files in space
            files = conn.execute("""
                SELECT name, file_type, file_path
                FROM datasets
                WHERE space_id = ?
                ORDER BY created_at DESC
                LIMIT ?
            """, [space_id, request.limit]).fetchall()
            
            print(f"[External API] Found {len(files)} files")
            
            total_count = conn.execute("""
                SELECT COUNT(*) FROM datasets WHERE space_id = ?
            """, [space_id]).fetchone()[0]
        
        # Read content from files
        documents = []
        for file_name, file_type, file_path in files:
            print(f"[External API] Reading file: {file_name}")
            content = await _read_file_content(file_path)
            if content:
                documents.append(DocumentContent(
                    file_name=file_name,
                    file_type=file_type,
                    content=content[:10000]  # Limit content size
                ))
        
        # Record API usage
        api_key_service.record_usage(key_info["id"], "data_query")
        
        print(f"[External API] Returning {len(documents)} documents")
        return DataQueryResponse(
            space_name=space[1],
            documents=documents,
            total_documents=total_count
        )
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        print(f"[External API] ERROR: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/space-info")
async def get_space_info(x_api_key: str = Header(..., alias="X-API-Key")):
    """Get information about the space linked to this API key"""
    key_info = await validate_api_key(x_api_key)
    
    space_id = key_info.get("space_id")
    if not space_id:
        return {"error": "API Key chưa gắn với space nào", "space": None}
    
    with get_db() as conn:
        space = conn.execute("""
            SELECT id, name, description, file_count, total_size_mb
            FROM document_spaces WHERE id = ?
        """, [space_id]).fetchone()
        
        if not space:
            return {"error": "Space không tồn tại", "space": None}
        
        return {
            "space": {
                "id": space[0],
                "name": space[1],
                "description": space[2],
                "file_count": space[3],
                "total_size_mb": space[4]
            }
        }


async def _read_file_content(file_path: str) -> Optional[str]:
    """Read content from various file types"""
    path = Path(file_path)
    
    if not path.exists():
        return None
    
    ext = path.suffix.lower()
    
    try:
        # Text files
        if ext in ['.txt', '.csv', '.json', '.md']:
            return path.read_text(encoding='utf-8', errors='ignore')[:10000]
        
        # PDF
        elif ext == '.pdf':
            import pdfplumber
            text = []
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages[:10]:
                    text.append(page.extract_text() or "")
            return "\n".join(text)[:10000]
        
        # Word docs
        elif ext == '.docx':
            from docx import Document
            doc = Document(file_path)
            text = "\n".join([p.text for p in doc.paragraphs])
            return text[:10000]
        
        # Excel
        elif ext in ['.xlsx', '.xls']:
            import pandas as pd
            df = pd.read_excel(file_path, nrows=100)
            return df.to_string()[:10000]
        
        # PowerPoint
        elif ext == '.pptx':
            from pptx import Presentation
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides[:10]:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)[:10000]
        
        # Parquet
        elif ext == '.parquet':
            import pandas as pd
            df = pd.read_parquet(file_path).head(100)
            return df.to_string()[:10000]
        
        else:
            return None
            
    except Exception as e:
        print(f"[External API] Read error {file_path}: {e}")
        return None


# =============================================================================
# INGEST ENDPOINT - Upload files via API Key
# =============================================================================

class IngestResponse(BaseModel):
    success: bool
    file_id: str
    file_name: str
    message: str


@router.post("/ingest", response_model=IngestResponse)
async def ingest_file(
    request: Request,
    x_api_key: str = Header(..., alias="X-API-Key")
):
    """
    Upload a file directly via API Key.
    
    Third parties can upload documents without accessing the Dashboard.
    Files are stored in the space linked to the API Key.
    
    Headers:
        X-API-Key: Your Locaith API Key (LcAi_xxx)
        Content-Type: multipart/form-data
    
    Body:
        file: The file to upload (PDF, DOCX, TXT, CSV, XLSX, etc.)
        name: (optional) Custom name for the file
        description: (optional) Description of the file
    """
    try:
        from fastapi import Form, UploadFile, File
        import uuid
        import os
        
        key_info = await validate_api_key(x_api_key)
        space_id = key_info.get("space_id")
        
        if not space_id:
            raise HTTPException(
                status_code=400,
                detail="API Key chưa được gắn với Document Space. Vui lòng tạo API Key từ trong Space."
            )
        
        # Parse multipart form
        form = await request.form()
        file = form.get("file")
        name = form.get("name", "")
        description = form.get("description", "Uploaded via External API")
        
        if not file:
            raise HTTPException(status_code=400, detail="Không tìm thấy file trong request")
        
        # Generate file ID and path
        file_id = str(uuid.uuid4())
        original_name = file.filename
        ext = os.path.splitext(original_name)[1].lower()
        
        # Validate file type
        allowed_types = ['.pdf', '.docx', '.txt', '.csv', '.xlsx', '.xls', '.json', '.md', '.pptx']
        if ext not in allowed_types:
            raise HTTPException(
                status_code=400, 
                detail=f"Loại file không được hỗ trợ. Cho phép: {', '.join(allowed_types)}"
            )
        
        # Save file
        from config import settings
        raw_dir = os.path.join(settings.DATA_DIR, "raw")
        os.makedirs(raw_dir, exist_ok=True)
        
        file_path = os.path.join(raw_dir, f"{file_id}{ext}")
        content = await file.read()
        
        with open(file_path, "wb") as f:
            f.write(content)
        
        file_size = len(content)
        display_name = name if name else original_name.replace(ext, "")
        
        # Use lakehouse service for 100% accuracy pipeline (includes AI OCR & Cleaning)
        from services.lakehouse import lakehouse_service
        result = await lakehouse_service.ingest_file_by_api_key(
            file_path=file_path,
            api_key_id=key_info["id"],
            user_id=key_info["user_id"],
            name=display_name,
            description=description
        )
        
        # Cleanup temp file if needed (ingest_file handles storage if configured)
        # However, external ingest uses settings.DATA_DIR/raw directly.
        # ingest_file will save it to datasets/ parquet.
        
        # Record usage
        api_key_service.record_usage(key_info["id"], "file_upload")
        
        return IngestResponse(
            success=True,
            file_id=result["id"],
            file_name=result["name"],
            message=f"File '{result['name']}' đã được upload và chuẩn hóa thành công!"
        )
        
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        print(f"[External API] Ingest error: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))


# =============================================================================
# QUERY ENDPOINT - Ask AI questions about your data
# =============================================================================

class QueryRequest(BaseModel):
    question: str
    max_documents: int = 5


class QueryResponse(BaseModel):
    answer: str
    sources: List[str]
    tokens_used: int


@router.post("/query", response_model=QueryResponse)
async def query_with_ai(
    request: QueryRequest,
    x_api_key: str = Header(..., alias="X-API-Key")
):
    """
    Ask AI questions about your documents.
    
    This endpoint retrieves relevant documents from your space and uses
    Locaith's AI (Gemini) to generate an answer based on the content.
    
    Headers:
        X-API-Key: Your Locaith API Key (LcAi_xxx)
    
    Body:
        question: Your question about the documents
        max_documents: (optional) Maximum documents to use as context (default: 5)
    """
    try:
        import google.generativeai as genai
        from config import settings
        
        key_info = await validate_api_key(x_api_key)
        space_id = key_info.get("space_id")
        
        if not space_id:
            raise HTTPException(
                status_code=400,
                detail="API Key chưa được gắn với Document Space"
            )
        
        # Get documents
        with get_db() as conn:
            space = conn.execute(
                "SELECT name FROM document_spaces WHERE id = ?", 
                [space_id]
            ).fetchone()
            
            if not space:
                raise HTTPException(status_code=404, detail="Space không tồn tại")
            
            files = conn.execute("""
                SELECT name, file_type, file_path
                FROM datasets WHERE space_id = ?
                ORDER BY created_at DESC LIMIT ?
            """, [space_id, request.max_documents]).fetchall()
        
        if not files:
            return QueryResponse(
                answer="Không có tài liệu nào trong space này. Vui lòng upload file trước.",
                sources=[],
                tokens_used=0
            )
        
        # Read content from files
        context_parts = []
        sources = []
        
        for file_name, file_type, file_path in files:
            content = await _read_file_content(file_path)
            if content:
                context_parts.append(f"=== {file_name}.{file_type} ===\n{content[:5000]}")
                sources.append(f"{file_name}.{file_type}")
        
        if not context_parts:
            return QueryResponse(
                answer="Không thể đọc nội dung file. Vui lòng kiểm tra định dạng file.",
                sources=[],
                tokens_used=0
            )
        
        # Generate answer with Gemini using new SDK
        context = "\n\n".join(context_parts)
        
        from google import genai
        client = genai.Client(api_key=settings.GEMINI_API_KEY)
        
        prompt = f"""Bạn là trợ lý AI chuyên phân tích tài liệu. Dựa trên nội dung sau, hãy trả lời câu hỏi một cách chính xác và súc tích.

TÀI LIỆU:
{context}

CÂU HỎI: {request.question}

Hãy trả lời bằng tiếng Việt, dựa trên thông tin trong tài liệu. Nếu không tìm thấy thông tin liên quan, hãy nói rõ."""

        response = client.models.generate_content(
            model='gemini-3-flash-preview',
            contents=prompt
        )
        answer = response.text
        
        # Estimate tokens (rough)
        tokens_used = len(prompt.split()) + len(answer.split())
        
        # Record usage
        api_key_service.record_usage(key_info["id"], "ai_query")
        
        return QueryResponse(
            answer=answer,
            sources=sources,
            tokens_used=tokens_used
        )
        
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        print(f"[External API] Query error: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))
